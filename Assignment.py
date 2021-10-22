#You will need python-pptx - You will be creating a ppt with 5 pages.
#You will need ImageMagick.
#You will need Wand - an ImageMagick binding for Python.

#Assignment - You need to write a python script that takes the files ( images + logo) from a given directory, 
#and outputs the resulting ppt file in the same directory. 
#There are 5 images and 1 logo file (a Nike logo) provided to you in this folder: (please download to your local)
#1. Please watermark the 5 image files with the Nike logo provided, using ImageMagick.
#(See sample file for placement, top left corner)
#2. Insert each watermarked image into a new page on the ppt. 
#3. Maintain the original aspect ratio of each image.
#Each page of the ppt should have a Heading, and Subheading, in a font of your choice. 

#Note - You will need to write the code in python. Please make sure to use a FOR loop, and good variables, so that your code is neat! 

#Sample output:
#A sample of the desired output is provided in the above folder. Please make sure your code outputs something similar. 

#Common Errors: 
#A. Logo placement  
#B. Distorting the aspect ratio of the logo/images 
#C. Not using a simple FOR loop  
#D. Too many variables, bad messy code  
#E. Not putting ImageMagick code in the .Py script 
#F. Hard-coding instead of using a loop G. ‘Cropping’ the images so it fits.


from pptx import Presentation
from pptx.util import Inches
from wand.image import Image
from wand.compat import nested

new_images=[]
images = [
  'image1.jpg',
  'image2.jpg',
  'image3.jpg',
  'image4.jpg',
  'image5.jpg'
]
 
logo = Image(filename='nike_black.png')

for i in images:
    logo.resize(840,300,filter='undefined')
    image = Image(filename=i)
    logo.transparentize(0)
    image.composite(logo,5,5)
    filename1='output_'+i
    image.save(filename='output_'+i)
    new_images.append(filename1)

prs = Presentation()


top = Inches(2.5)
left = Inches(1)
height = Inches(4.5)


class MySlide:
    def __init__(self, data):
       
        self.layout = prs.slide_layouts[data[3]]
  
        self.slide= prs.slides.add_slide(self.layout)
   
        self.heading= self.slide.shapes.title
        self.heading.text= data[0]
       
        self.sub_heading=self.slide.placeholders[1]
        self.sub_heading.text=data[1]
        
        if data[2] != "":          
            self.slide.shapes.add_picture(data[2], left, top, height=height)



slides = [
    ["Sample Title 1",       
     "Sample Subtitle 1",
     new_images[0],
     1],
    ["Sample Title 2",      
     "Sample Subtitle 2",
     new_images[1],
     1],
    ["Sample Title 3",       
     "Sample Subtitle 3",
     new_images[2],
     1],
     ["Sample Title 4",       
     "Sample Subtitle 4",
     new_images[3],
     1],
     ["Sample Title 5",       
     "Sample Subtitle 5",
     new_images[4],
     1]
]


for each_slide in slides:
    MySlide(each_slide)


prs.save('sample_ppt.pptx')

    


